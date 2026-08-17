import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_classroom_ppt():
    prs = Presentation()
    
    # Set Slide Size to 16:9 widescreen
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Styling variables
    font_header = 'Arial'
    font_body = 'Calibri'
    color_primary = RGBColor(10, 37, 64)   # Deep Blue
    color_accent = RGBColor(0, 128, 128)   # Teal
    color_charcoal = RGBColor(51, 51, 51)  # Charcoal
    color_bg = RGBColor(248, 250, 252)     # Slate-50 (Premium clean background)
    
    # Image Paths setup
    base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    assets_dir = os.path.join(base_dir, "assets")
    
    img_landing_dark = os.path.join(assets_dir, "landing_page_dark.png")
    img_python_trace = os.path.join(assets_dir, "python_trace.png")
    img_visualize_dark = os.path.join(assets_dir, "visualize_dark.png")
    img_visualize_light = os.path.join(assets_dir, "visualize_light.png")
    img_landing_light = os.path.join(assets_dir, "landing_page.png")
    
    slide_layout = prs.slide_layouts[6] # Blank slide
    
    # ─── SLIDE 1: Welcome & Hook ───
    slide1 = prs.slides.add_slide(slide_layout)
    apply_slide_background(slide1, color_bg)
    
    # Left Text Column
    tb_text = slide1.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(6.0), Inches(5.0))
    tf = tb_text.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    run_title = p1.add_run()
    run_title.text = "Why Do Coding Labs Feel Like a Nightmare?"
    run_title.font.name = font_header
    run_title.font.size = Pt(36)
    run_title.font.bold = True
    run_title.font.color.rgb = color_primary
    
    p2 = tf.add_paragraph()
    p2.space_before = Pt(14)
    run_sub = p2.add_run()
    run_sub.text = "CodeCanvas (LPU CodeViz)"
    run_sub.font.name = font_header
    run_sub.font.size = Pt(18)
    run_sub.font.bold = True
    run_sub.font.color.rgb = color_accent
    
    p3 = tf.add_paragraph()
    p3.space_before = Pt(24)
    run_meta = p3.add_run()
    run_meta.text = (
        "• You write syntax, compile, and get stuck on logical bugs.\n"
        "• Lab instructors are busy helping 15 other students.\n"
        "• Debugging terminals are complex for beginners."
    )
    run_meta.font.name = font_body
    run_meta.font.size = Pt(16)
    run_meta.font.color.rgb = color_charcoal
    
    # Right Image Column
    if os.path.exists(img_landing_dark):
        slide1.shapes.add_picture(img_landing_dark, Inches(7.2), Inches(1.2), width=Inches(5.3), height=Inches(5.0))
    
    slide1.notes_slide.notes_text_frame.text = (
        "Hey everyone! Let's be honest for a second. How many of you have sat in a programming lab, "
        "written code that you thought was 100% correct, only for it to return a completely wrong output? "
        "Or worse, you get a 'Segmentation Fault' or 'Index Out of Range' and have absolutely no idea "
        "which line actually broke? You raise your hand, but the lab teacher is already surrounded by "
        "15 other students. You end up wasting half the lab session waiting for help, or coding late at night "
        "with zero support. That is the coding lab gap, and today, we're changing that."
    )
    
    # ─── SLIDE 2: The Core Problem ───
    slide2 = prs.slides.add_slide(slide_layout)
    apply_slide_background(slide2, color_bg)
    
    tb_text = slide2.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(6.0), Inches(5.0))
    tf = tb_text.text_frame
    tf.word_wrap = True
    
    p_hdr = tf.paragraphs[0]
    run_hdr = p_hdr.add_run()
    run_hdr.text = "The Core Problem: Memory is Invisible"
    run_hdr.font.name = font_header
    run_hdr.font.size = Pt(36)
    run_hdr.font.bold = True
    run_hdr.font.color.rgb = color_primary
    
    p_body = tf.add_paragraph()
    p_body.space_before = Pt(24)
    run_body = p_body.add_run()
    run_body.text = (
        "• Coding Blind: Standard code compilers run your scripts, but keep memory layouts invisible.\n"
        "• Theoretical Lectures: Pointers, stacks, and node links are explained statically on blackboards.\n"
        "• Guesswork Debugging: Beginners lack a clear mental model of memory, making debugging a guessing game."
    )
    run_body.font.name = font_body
    run_body.font.size = Pt(16)
    run_body.font.color.rgb = color_charcoal
    
    if os.path.exists(img_python_trace):
        slide2.shapes.add_picture(img_python_trace, Inches(7.2), Inches(1.2), width=Inches(5.3), height=Inches(5.0))
        
    slide2.notes_slide.notes_text_frame.text = (
        "The reason programming feels so intimidating when you start is that computer memory is invisible. "
        "When a professor explains pointers referencing memory, nodes swapping in a linked list, or stacks popping "
        "in a recursion loop, they write it statically on a blackboard. But when you sit in front of VS Code, "
        "you can't see the memory stack. You are essentially coding blind. You compile, it fails, and you guess. "
        "What if you could actually see your variables mutating, your arrays sorting, and your linked lists linking, "
        "step-by-step in real-time?"
    )
    
    # ─── SLIDE 3: The Solution ───
    slide3 = prs.slides.add_slide(slide_layout)
    apply_slide_background(slide3, color_bg)
    
    tb_text = slide3.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(6.0), Inches(5.0))
    tf = tb_text.text_frame
    tf.word_wrap = True
    
    p_hdr = tf.paragraphs[0]
    run_hdr = p_hdr.add_run()
    run_hdr.text = "The Solution: CodeCanvas (LPU CodeViz)"
    run_hdr.font.name = font_header
    run_hdr.font.size = Pt(36)
    run_hdr.font.bold = True
    run_hdr.font.color.rgb = color_primary
    
    p_body = tf.add_paragraph()
    p_body.space_before = Pt(24)
    run_body = p_body.add_run()
    run_body.text = (
        "• Monaco Coding Editor: Paste and edit arbitrary C, Python, or SQL scripts.\n"
        "• 10 Animated Canvases: Play step-by-step memory animations (Arrays, Stacks, Trees, Graphs, SQL).\n"
        "• Variables Board: Observe variables mutating scopes in real-time.\n"
        "• AI Tutor Companion: Slides out to offer hint instructions, helping you debug logic errors."
    )
    run_body.font.name = font_body
    run_body.font.size = Pt(16)
    run_body.font.color.rgb = color_charcoal
    
    if os.path.exists(img_visualize_dark):
        slide3.shapes.add_picture(img_visualize_dark, Inches(7.2), Inches(1.2), width=Inches(5.3), height=Inches(5.0))
        
    slide3.notes_slide.notes_text_frame.text = (
        "This is CodeCanvas, also known as LPU CodeViz. It's a web-based coding environment built specifically for us. "
        "You write your C, C++, Python, or SQL code in a built-in editor, click 'Visualize', and the system generates "
        "a step-by-step animated playback of your program. You can watch your array bars swap during Bubble Sort, "
        "watch your recursion call stack grow and shrink, and see variable scopes update instantly. Plus, there's a "
        "built-in AI Tutor sidebar that acts like a personal teacher, checking your specific code and guiding you to "
        "fix bugs without giving away answers."
    )
    
    # ─── SLIDE 4: Student Benefits ───
    slide4 = prs.slides.add_slide(slide_layout)
    apply_slide_background(slide4, color_bg)
    
    tb_text = slide4.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(6.0), Inches(5.0))
    tf = tb_text.text_frame
    tf.word_wrap = True
    
    p_hdr = tf.paragraphs[0]
    run_hdr = p_hdr.add_run()
    run_hdr.text = "What's In It For You?"
    run_hdr.font.name = font_header
    run_hdr.font.size = Pt(36)
    run_hdr.font.bold = True
    run_hdr.font.color.rgb = color_primary
    
    p_body = tf.add_paragraph()
    p_body.space_before = Pt(24)
    run_body = p_body.add_run()
    run_body.text = (
        "• Understand Instantly: Dynamic visual debuggers establish clear, correct mental models.\n"
        "• Debug 24/7: Get immediate guidance late at night—no more waiting for lab office hours.\n"
        "• Gamified Tracking: Build coding streaks, earn XP, and unlock cryptographic certificates.\n"
        "• Syllabus Aligned: Matches standard LPU engineering assignments and course codes."
    )
    run_body.font.name = font_body
    run_body.font.size = Pt(16)
    run_body.font.color.rgb = color_charcoal
    
    if os.path.exists(img_visualize_light):
        slide4.shapes.add_picture(img_visualize_light, Inches(7.2), Inches(1.2), width=Inches(5.3), height=Inches(5.0))
        
    slide4.notes_slide.notes_text_frame.text = (
        "So how does this help you pass your exams and build real coding skills? First, it gives you an instant "
        "mental model of memory, saving you hours of frustration. Second, you can debug your code late at night—"
        "it's like having a coding tutor in your browser 24/7. Third, we've pre-loaded it with standard syllabus exercises. "
        "Finally, it's gamified: you earn XP, build coding streaks, and once you complete a topic, you get a "
        "cryptographically verified certificate shareable on LinkedIn to stand out to recruiters."
    )
    
    # ─── SLIDE 5: CTA ───
    slide5 = prs.slides.add_slide(slide_layout)
    apply_slide_background(slide5, color_bg)
    
    tb_text = slide5.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(6.0), Inches(5.0))
    tf = tb_text.text_frame
    tf.word_wrap = True
    
    p_hdr = tf.paragraphs[0]
    run_hdr = p_hdr.add_run()
    run_hdr.text = "Get Registered in 30 Seconds"
    run_hdr.font.name = font_header
    run_hdr.font.size = Pt(36)
    run_hdr.font.bold = True
    run_hdr.font.color.rgb = color_primary
    
    p_body = tf.add_paragraph()
    p_body.space_before = Pt(24)
    run_body = p_body.add_run()
    run_body.text = (
        "• Visit: codecanvas.edu/register\n"
        "• Register: One-click sign-up with Google OAuth or Email.\n"
        "• Join Classroom: Input invite code: [CLASS_INVITE_CODE]\n"
        "• Start Coding: Visualize your code and start building your streaks today!"
    )
    run_body.font.name = font_body
    run_body.font.size = Pt(16)
    run_body.font.color.rgb = color_charcoal
    
    if os.path.exists(img_landing_light):
        slide5.shapes.add_picture(img_landing_light, Inches(7.2), Inches(1.2), width=Inches(5.3), height=Inches(5.0))
        
    slide5.notes_slide.notes_text_frame.text = (
        "Signing up is incredibly fast. Pull out your phone right now, scan the QR code on the screen, "
        "or open codecanvas.edu in your browser. Click 'Sign Up', log in instantly with your Google account, "
        "and select the 'Student' role. To join our lab section and track your outcomes, simply enter the "
        "classroom invite code shown here. It takes less than 30 seconds. Register now, start visualizing "
        "your code, and let's make coding labs a breeze this semester. Thanks everyone!"
    )
    
    # Save the presentation
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "classroom_pitch.pptx")
    try:
        prs.save(out_path)
        print(f"Aesthetic Classroom Pitch PPT generated successfully at: {out_path}")
    except PermissionError:
        fallback_path = os.path.join(out_dir, "classroom_pitch_updated.pptx")
        prs.save(fallback_path)
        print(f"File locked by another application. Saved instead to fallback: {fallback_path}")

if __name__ == "__main__":
    create_classroom_ppt()
